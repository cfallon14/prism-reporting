# Python utility
import argparse
import os
import sys
import shutil
from abc import ABC, abstractmethod

# Data processing
import pandas as pd
import json

# Reporting
import xlsxwriter

from jinja2 import Environment, FileSystemLoader
from weasyprint import HTML

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def main():
    #parser = argparse.ArgumentParser(description='PrismReporting suite.')
    #parser.add_argument('--global_settings', dest='setting_file', help='Global setting file')

    #args = parser.parse_args(['--global_settings'])
    #print(args.setting_file)

    file = sys.argv[0]
    prism_dir = os.path.dirname(file)

    cwd = os.getcwd()

    # Open up the global settings file
    try:
        with open(os.path.join(cwd, 'settings.json')) as global_settings_file:
            settings = json.load(global_settings_file)
    except (OSError, IOError) as e:
        print('Error with settings.json file. Please add or modify settings.json in current directory.')
        raise

    # PDF Settings are defined
    if 'pdf_settings' in settings:
        pdf_settings = settings['pdf_settings']

        report_template = pdf_settings['templates']
        report_name = settings['report_name']

        dst = ''.join(['./', report_name,'/pdf'])

        if not os.path.isdir(dst):
            shutil.copytree(''.join([prism_dir, '/templates/pdf/template_',
                                 report_template]), dst)

        # Create result directory in report folder
        if not os.path.exists(''.join(['./',report_name, '/result'])):
            os.makedirs(''.join(['./',report_name, '/result']))

        pdf_report = PrismPDF(settings)
        pdf_report.run_report()

    # Powerpoint settings are defined
    if 'ppt_settings' in settings:
        ppt_settings = settings['ppt_settings']

        if ppt_settings:
            ppt_report = PrismPPT(settings)
            ppt_report.run_report()

    # Excel settings are defined
    if 'xl_settings' in settings:
        xl_settings = settings['xl_settings']

        if xl_settings:
            xl_report = PrismXL(settings)
            xl_report.run_report()


class PrismReport(ABC):
    """ Abstract class for reporting
    """
    @abstractmethod
    def __init__(self, settings):
        self._project_path = settings['project_path']
        self._report_name = settings['report_name']
        self._data_file = settings['data_file']
        self._template_vars = settings['template_vars']
        self._output_filename = settings['output_filename']
        self._styles = []
        self._template = ''

    @abstractmethod
    def init_report(self, df):
        pass

    @abstractmethod
    def add_report_page(self, df):
        pass


    @abstractmethod
    def set_charts(self, df):
        pass

    @abstractmethod
    def set_tables(self, df):
        pass

    @abstractmethod
    def set_styles(self):
        pass

    @abstractmethod
    def run_report(self):
        pass

    def get_project_path(self):
        return self._project_path

    def get_data_file(self):
        return self._data_file

    def get_template_vars(self):
        return self._template_vars

    def get_output_filename(self):
        return self._output_filename

    def get_styles(self):
        return self._styles

    def get_template(self):
        return self._template

    class PrismTable(ABC):
        def create_table(self):
            pass


class PrismPDF(PrismReport):

    def __init__(self, settings):
        PrismReport.__init__(self, settings)
        self._pdf_settings = settings['pdf_settings']

    def init_report(self, df):
        pass

    def add_report_page(self, df):
        pass

    def set_charts(self, df):
        self._template_vars['charts'] = {}

    def set_tables(self, df):
        self._template_vars['tables'] = [value.to_html(index=False) for key, value in df.items()]

    def set_styles(self):
        self._styles = [''.join(['./pdf/styles/', s]) for s in os.listdir('./pdf/styles/')]

    def run_report(self):
        # Report path
        report_path = ''.join([self._project_path, '/', self._report_name])

        # Change to directory of report
        os.chdir(report_path)

        # Output will be pdf
        output_file = ''.join([self._output_filename,'.pdf'])

        template_file = self._pdf_settings['templates']

        # Import data from excel file into Pandas dataframe
        df = pd.read_excel(''.join(['./data/', self._data_file]), sheetname=None)

        env = Environment(loader=FileSystemLoader(report_path))
        template = env.get_template(''.join(['./pdf/', template_file,'.html']))

        # Convert tabular data into html tables
        self.set_tables(df)

        html_out = template.render(self._template_vars)

        # Get all the stylesheets
        self.set_styles()

        # Base url needs to be project folder name
        base_url = self._project_path.split('/')[-1]
        HTML(string=html_out, base_url=base_url).write_pdf(''.join(['./result/', output_file]),
                                                           stylesheets=self._styles)


class PrismXL(PrismReport):

    def __init__(self, settings):
        PrismReport.__init__(self, settings)

    def set_charts(self, df):
        pass

    def set_tables(self, df):
        pass

    def set_styles(self):
        pass

    def run_report(self):
        workbook = xlsxwriter.Workbook('hello.xlsx')
        worksheet = workbook.add_worksheet()

        worksheet.write('A1', 'Hello world')

        workbook.close()


class PrismPPT(PrismReport):

    def __init__(self, settings):
        PrismReport.__init__(self,settings)

    def set_charts(self, df):
        pass

    def set_tables(self, df):
        pass

    def set_styles(self):
        pass

    def run_report(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = ChartData()
        chart_data.categories = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10','11','12','13', '14','15','16']
        chart_data.add_series('West',    (32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 14))
        chart_data.add_series('East',    (24.3, 30.6, 20.2, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 13))
        chart_data.add_series('Midwest', (20.4, 18.3, 26.2, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 32.2, 28.4, 34.7, 12))

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.include_in_layout = False

        prs.save('chart-01.pptx')


if __name__ == '__main__':
    main()